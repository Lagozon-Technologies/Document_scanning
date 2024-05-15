import os
import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
from langchain.callbacks import get_openai_callback
from PIL import Image
import pandas as pd
import json
from docx import Document
from pptx import Presentation
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.llms import OpenAI
from langchain_community.callbacks import get_openai_callback
import toml


api_key = st.secrets["API_KEY"]

img = Image.open(r"imagess.png")
st.set_page_config(page_title="InsightVault: Document Generation AI", page_icon=img)
col1, col2 = st.columns([1, 4])  


with col1:
    st.image("img.jpg", width=90)


with col2:
    st.title("Ask Your Documents 📄")

file_type = st.selectbox("Select file type", ["PDF", "DOCX", "PPTX"])
accepted_types = {"PDF": "pdf", "DOCX": "docx", "PPTX": "pptx"}
file_extension = accepted_types[file_type]

uploaded_files = st.file_uploader(f"Upload your {file_type} files", type=file_extension, accept_multiple_files=True)

if uploaded_files:
    all_texts = []
    for uploaded_file in uploaded_files:
        if file_type == "PDF":
            pdf_reader = PdfReader(uploaded_file)
            text = "".join([page.extract_text() or "" for page in pdf_reader.pages])
            if text:
                all_texts.append(text)
        elif file_type == "DOCX":
            doc = Document(uploaded_file)
            doc_text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            all_texts.append(doc_text)
        elif file_type == "PPTX":
            ppt = Presentation(uploaded_file)
            ppt_text = '\n'.join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
            all_texts.append(ppt_text)

    if all_texts:
        text_splitter = CharacterTextSplitter(separator="\n", chunk_size=1000, chunk_overlap=200, length_function=len)
        chunks = []
        for text in all_texts:
            chunks.extend(text_splitter.split_text(text))

        embeddings = OpenAIEmbeddings(api_key=api_key)
        knowledge_base = FAISS.from_texts(chunks, embeddings)

        query = st.text_input("Ask your Question about your documents")
        if query:
            docs = knowledge_base.similarity_search(query)
            llm = OpenAI(api_key=api_key)
            chain = load_qa_chain(llm, chain_type="stuff")
            response = chain.run(input_documents=docs, question=query)
            st.success(response)
